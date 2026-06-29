#!/usr/bin/env python3
"""Systematic Factor Income Fund — strategy & pipeline presentation (.pptx).

Comprehensive walkthrough for an internal / leadership audience. Plain-language
explanations of every term, prominent rendered formulas, explanatory graphics,
and a sober, factual tone. Editable and regenerable.

    <scratchpad>/deckvenv/bin/python presentation/build_deck.py

Output: presentation/Systematic-Factor-Income-Fund.pptx  (+ assets/*.png)

Benchmark figures (S&P 500, BXMD) are computed from index history over the
Sep 2021 - Jun 2026 window; strategy figures are from the walk-forward backtest.
"""
from __future__ import annotations
import os

import math
import textwrap

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import (FancyBboxPatch, FancyArrowPatch, Rectangle, Circle,
                                Polygon, RegularPolygon, Arc, Wedge)
from matplotlib.font_manager import FontProperties

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
os.makedirs(ASSETS, exist_ok=True)

# --------------------------------------------------------------------------- #
INK       = RGBColor(0x0E, 0x1B, 0x2A)
INK_PANEL = RGBColor(0x16, 0x29, 0x3C)
SLATE     = RGBColor(0x2A, 0x47, 0x60)
ACCENT    = RGBColor(0x14, 0xB8, 0xA6)
ACCENT_DK = RGBColor(0x0E, 0x8C, 0x7E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0x1F, 0x2A, 0x37)
MUTED     = RGBColor(0x6B, 0x77, 0x87)
RULE      = RGBColor(0xE3, 0xE8, 0xEE)
PANEL     = RGBColor(0xF4, 0xF7, 0xF9)
PANEL_BD  = RGBColor(0xDD, 0xE4, 0xEA)
LIGHT     = RGBColor(0xC7, 0xD2, 0xDD)
FONT = "Avenir Next"

INK_H, ACC_H, ACCDK_H = "#0E1B2A", "#14B8A6", "#0E8C7E"
MUT_H, GRY_H, TEXT_H = "#6B7787", "#9AA7B4", "#1F2A37"
PANEL_H, BD_H, LIGHT_H = "#F4F7F9", "#DDE4EA", "#C7D2DD"
_BOLD = FontProperties(family="Avenir Next", weight="bold")
_REG = FontProperties(family="Avenir Next")

matplotlib.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Avenir Next", "Helvetica Neue", "Arial", "DejaVu Sans"],
    "mathtext.fontset": "cm",
    "axes.edgecolor": MUT_H, "axes.labelcolor": TEXT_H,
    "text.color": TEXT_H, "xtick.color": TEXT_H, "ytick.color": TEXT_H,
    "axes.linewidth": 0.9, "figure.dpi": 200,
})


# =========================================================================== #
# Asset generation
# =========================================================================== #
def formula(tex, name, fontsize=40, color=INK_H):
    path = os.path.join(ASSETS, name)
    fig = plt.figure(figsize=(0.1, 0.1))
    fig.text(0.0, 0.0, f"${tex}$", fontsize=fontsize, color=color)
    fig.savefig(path, dpi=300, bbox_inches="tight", pad_inches=0.08, transparent=True)
    plt.close(fig)
    return path


def fig_bell():
    path = os.path.join(ASSETS, "bell.png")
    x = np.linspace(-3.6, 3.6, 500)
    y = np.exp(-x**2 / 2) / np.sqrt(2 * np.pi)
    fig, ax = plt.subplots(figsize=(7.2, 3.2))
    ax.plot(x, y, color=INK_H, lw=2)
    ax.fill_between(x, y, where=(x >= -1) & (x <= 1), color=ACC_H, alpha=0.18)
    ax.axvline(0, color=MUT_H, lw=1, ls=(0, (4, 3)))
    ax.axvline(1.2, color=ACCDK_H, lw=2)
    ax.scatter([1.2], [np.exp(-1.2**2 / 2) / np.sqrt(2 * np.pi)], color=ACCDK_H, zorder=5, s=30)
    ax.annotate("average company\n(score 0)", xy=(0, 0.02), xytext=(-3.0, 0.28),
                fontsize=11, color=MUT_H, ha="left",
                arrowprops=dict(arrowstyle="->", color=MUT_H, lw=1))
    ax.annotate("this company\nscore +1.2", xy=(1.2, 0.194), xytext=(1.75, 0.30),
                fontsize=11, color=ACCDK_H, ha="left", weight="bold",
                arrowprops=dict(arrowstyle="->", color=ACCDK_H, lw=1.2))
    ax.set_xlabel("standard deviations from the average company", fontsize=11)
    ax.set_yticks([]); ax.set_ylim(0, 0.46); ax.set_xticks(range(-3, 4))
    for sp in ("top", "right", "left"):
        ax.spines[sp].set_visible(False)
    fig.tight_layout(); fig.savefig(path, dpi=200, bbox_inches="tight", transparent=True)
    plt.close(fig); return path


def fig_payoff():
    path = os.path.join(ASSETS, "payoff.png")
    S0, K, prem = 100.0, 107.0, 2.0
    p = np.linspace(82, 128, 400)
    stock = p - S0
    cc = np.minimum(p, K) - S0 + prem
    fig, ax = plt.subplots(figsize=(7.4, 3.7))
    ax.axhline(0, color=MUT_H, lw=0.9)
    ax.plot(p, stock, color=GRY_H, lw=1.8, ls=(0, (5, 3)), label="Hold the stock only")
    ax.plot(p, cc, color=ACCDK_H, lw=2.6, label="Stock plus covered call")
    ax.axvline(K, color=MUT_H, lw=1, ls=(0, (2, 3)))
    ax.axhline(prem, color=ACC_H, lw=1, ls=(0, (2, 3)))
    ax.annotate("premium kept\n(+$2 / share)", xy=(86, prem), xytext=(84, 9),
                fontsize=10.5, color=ACCDK_H, ha="left",
                arrowprops=dict(arrowstyle="->", color=ACCDK_H, lw=1))
    ax.annotate("gains capped\nabove the strike", xy=(118, K - S0 + prem),
                xytext=(108.5, -6.5), fontsize=10.5, color=MUT_H, ha="left",
                arrowprops=dict(arrowstyle="->", color=MUT_H, lw=1))
    ax.text(K + 0.4, -13.5, "strike $107", fontsize=10, color=MUT_H)
    ax.set_xlabel("stock price at the option's expiration ($)", fontsize=11)
    ax.set_ylabel("profit / loss per share ($)", fontsize=11)
    ax.legend(loc="upper left", frameon=False, fontsize=10.5)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    fig.tight_layout(); fig.savefig(path, dpi=200, bbox_inches="tight", transparent=True)
    plt.close(fig); return path


def fig_risk_return():
    path = os.path.join(ASSETS, "risk_return.png")
    # label, vol%, ret%, color, highlight, dx, dy, ha
    pts = [
        ("S&P 500", 16.0, 13.3, GRY_H, False, 0.5, 0.25, "left"),
        ("S&P 500 BuyWrite", 12.0, 10.2, MUT_H, False, 0.0, -1.9, "center"),
        ("Strategy, stocks only", 20.0, 10.9, "#566374", False, 0.0, -1.9, "center"),
        ("Strategy plus income overlay", 15.5, 19.6, ACCDK_H, True, 0.45, 0.95, "left"),
    ]
    fig, ax = plt.subplots(figsize=(7.5, 4.4))
    for label, vol, ret, col, hi, dx, dy, ha in pts:
        ax.scatter([vol], [ret], s=250 if hi else 150, color=col,
                   edgecolor="white", linewidth=1.5, zorder=5)
        ax.annotate(label, (vol, ret), xytext=(vol + dx, ret + dy),
                    fontsize=11, color=col, weight="bold" if hi else "normal", ha=ha)
    ax.annotate("", xy=(15.0, 21.0), xytext=(19.8, 12.4),
                arrowprops=dict(arrowstyle="->", color=ACC_H, lw=1.4, alpha=0.5))
    ax.text(16.2, 21.3, "more return,\nless risk", fontsize=10.5, color=ACCDK_H, ha="left")
    ax.set_xlabel("risk:  annual volatility (%)", fontsize=11)
    ax.set_ylabel("annual return (%)", fontsize=11)
    ax.set_xlim(9, 24); ax.set_ylim(8, 23)
    ax.grid(True, color="#EEF2F5", lw=0.8); ax.set_axisbelow(True)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    fig.tight_layout(); fig.savefig(path, dpi=200, bbox_inches="tight", transparent=True)
    plt.close(fig); return path


def fig_vrp():
    path = os.path.join(ASSETS, "vrp.png")
    fig, ax = plt.subplots(figsize=(6.6, 3.5))
    bars = ax.bar(["Implied volatility\n(what we sell at)", "Realized volatility\n(what occurred)"],
                  [38.4, 33.2], color=[ACCDK_H, GRY_H], width=0.55)
    for b, v in zip(bars, [38.4, 33.2]):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.6, f"{v:.1f}%",
                ha="center", fontsize=12, weight="bold", color=INK_H)
    ax.annotate("", xy=(1, 33.6), xytext=(0, 38.0),
                arrowprops=dict(arrowstyle="<->", color=MUT_H, lw=1.2))
    ax.text(0.5, 41.0, "+5.2 points = the premium we harvest",
            ha="center", fontsize=11, color=ACCDK_H, weight="bold")
    ax.set_ylim(0, 46); ax.set_ylabel("annualized volatility (%)", fontsize=11)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    fig.tight_layout(); fig.savefig(path, dpi=200, bbox_inches="tight", transparent=True)
    plt.close(fig); return path


def _rbox(ax, x, y, w, h, fc=PANEL_H, ec=BD_H, lw=1.2, r=0.1):
    ax.add_patch(FancyBboxPatch((x, y), w, h,
                 boxstyle=f"round,pad=0,rounding_size={r}", fc=fc, ec=ec, lw=lw))


def _arc_arrow(ax, cx, cy, r, t1, t2, color, lw=2.2):
    ax.add_patch(Arc((cx, cy), 2 * r, 2 * r, theta1=t1, theta2=t2, ec=color, lw=lw))
    a = math.radians(t2)
    px, py = cx + r * math.cos(a), cy + r * math.sin(a)
    tx, ty = -math.sin(a), math.cos(a)
    nx, ny = math.cos(a), math.sin(a)
    sz = r * 0.6
    ax.add_patch(Polygon([(px + tx * sz, py + ty * sz),
                          (px + nx * sz * 0.5, py + ny * sz * 0.5),
                          (px - nx * sz * 0.5, py - ny * sz * 0.5)], closed=True, fc=color, ec="none"))


def _icon(ax, name, cx, cy, s, color=ACCDK_H):
    lw = 2.1
    if name == "shield":
        pts = [(cx, cy + s), (cx + 0.8*s, cy + 0.45*s), (cx + 0.8*s, cy - 0.25*s),
               (cx, cy - s), (cx - 0.8*s, cy - 0.25*s), (cx - 0.8*s, cy + 0.45*s)]
        ax.add_patch(Polygon(pts, closed=True, fill=False, ec=color, lw=lw, joinstyle="round"))
        ax.plot([cx - 0.32*s, cx - 0.05*s, cx + 0.4*s], [cy - 0.02*s, cy - 0.32*s, cy + 0.35*s],
                color=color, lw=lw, solid_capstyle="round")
    elif name == "stop":
        ax.add_patch(RegularPolygon((cx, cy), numVertices=8, radius=s, orientation=math.pi/8,
                     fill=False, ec=color, lw=lw))
        ax.plot([cx - 0.45*s, cx + 0.45*s], [cy, cy], color=color, lw=lw, solid_capstyle="round")
    elif name == "doc":
        _rbox(ax, cx - 0.55*s, cy - 0.8*s, 1.1*s, 1.6*s, fc="white", ec=color, lw=lw, r=0.04)
        for yy in (0.42, 0.05, -0.32):
            ax.plot([cx - 0.32*s, cx + 0.32*s], [cy + yy*s, cy + yy*s], color=color, lw=lw*0.8)
    elif name == "clock":
        ax.add_patch(Circle((cx, cy), s, fill=False, ec=color, lw=lw))
        ax.plot([cx, cx], [cy, cy + 0.6*s], color=color, lw=lw, solid_capstyle="round")
        ax.plot([cx, cx + 0.42*s], [cy, cy + 0.05*s], color=color, lw=lw, solid_capstyle="round")
    elif name == "check":
        ax.add_patch(Circle((cx, cy), s, fill=False, ec=color, lw=lw))
        ax.plot([cx - 0.42*s, cx - 0.1*s, cx + 0.45*s], [cy, cy - 0.35*s, cy + 0.42*s],
                color=color, lw=lw, solid_capstyle="round")
    elif name == "funnel":
        ax.add_patch(Polygon([(cx - 0.7*s, cy + 0.7*s), (cx + 0.7*s, cy + 0.7*s),
                              (cx + 0.16*s, cy - 0.1*s), (cx + 0.16*s, cy - 0.7*s),
                              (cx - 0.16*s, cy - 0.7*s), (cx - 0.16*s, cy - 0.1*s)],
                     closed=True, fill=False, ec=color, lw=lw, joinstyle="round"))
    elif name == "split":
        ax.plot([cx, cx], [cy - 0.7*s, cy], color=color, lw=lw, solid_capstyle="round")
        ax.plot([cx, cx - 0.6*s], [cy, cy + 0.55*s], color=color, lw=lw, solid_capstyle="round")
        ax.plot([cx, cx + 0.6*s], [cy, cy + 0.55*s], color=color, lw=lw, solid_capstyle="round")
        for ex in (-0.6, 0.6):
            ax.add_patch(Circle((cx + ex*s, cy + 0.55*s), 0.1*s, fc=color, ec="none"))
    elif name == "calendar":
        _rbox(ax, cx - 0.7*s, cy - 0.6*s, 1.4*s, 1.2*s, fc="white", ec=color, lw=lw, r=0.05)
        ax.plot([cx - 0.7*s, cx + 0.7*s], [cy + 0.26*s, cy + 0.26*s], color=color, lw=lw)
        for hx in (-0.35, 0.35):
            ax.plot([cx + hx*s, cx + hx*s], [cy + 0.45*s, cy + 0.78*s], color=color, lw=lw, solid_capstyle="round")
        for dx in (-0.32, 0, 0.32):
            for dy in (-0.05, -0.32):
                ax.add_patch(Circle((cx + dx*s, cy + dy*s), 0.07*s, fc=color, ec="none"))
    elif name == "percent":
        ax.add_patch(Circle((cx - 0.3*s, cy + 0.32*s), 0.2*s, fill=False, ec=color, lw=lw))
        ax.add_patch(Circle((cx + 0.3*s, cy - 0.32*s), 0.2*s, fill=False, ec=color, lw=lw))
        ax.plot([cx - 0.5*s, cx + 0.5*s], [cy - 0.55*s, cy + 0.55*s], color=color, lw=lw, solid_capstyle="round")
    elif name == "chartline":
        xs = [cx - 0.6*s, cx - 0.2*s, cx + 0.15*s, cx + 0.6*s]
        ys = [cy - 0.4*s, cy + 0.15*s, cy - 0.12*s, cy + 0.5*s]
        ax.plot(xs, ys, color=color, lw=lw, solid_capstyle="round")
        for x, y in zip(xs, ys):
            ax.add_patch(Circle((x, y), 0.08*s, fc=color, ec="none"))
    elif name == "refresh":
        _arc_arrow(ax, cx, cy, 0.78*s, 50, 300, color, lw)
    elif name == "bars":
        for i, hh in enumerate([0.5, 0.8, 1.15]):
            ax.add_patch(Rectangle((cx - 0.5*s + i*0.36*s, cy - 0.55*s), 0.24*s, hh*s, fc=color, ec="none"))


LOGO_DIR = os.path.join(HERE, "logos")
ALPACA_LOGO = os.path.join(LOGO_DIR, "alpaca.png")
SEC_LOGO = os.path.join(LOGO_DIR, "sec.png")


def _place_logo(ax, path, cx, cy, box):
    """Place a logo image centered at (cx, cy), fit within a 2*box square, aspect preserved."""
    img = plt.imread(path)
    ih, iw = img.shape[0], img.shape[1]
    ar = iw / ih
    if ar >= 1:
        w = 2 * box; h = w / ar
    else:
        h = 2 * box; w = h * ar
    ax.imshow(img, extent=[cx - w/2, cx + w/2, cy - h/2, cy + h/2],
              zorder=6, aspect="auto", interpolation="antialiased")


def _icon_cards(path, items, cols=2, figw=12.0, figh=4.5, wrap=34):
    rows = -(-len(items) // cols)
    fig, ax = plt.subplots(figsize=(figw, figh))
    ax.set_xlim(0, figw); ax.set_ylim(0, figh); ax.axis("off")
    pad = 0.28
    cw = (figw - pad * (cols + 1)) / cols
    ch = (figh - pad * (rows + 1)) / rows
    for i, (icon, title, desc) in enumerate(items):
        r = i // cols; c = i % cols
        x = pad + c * (cw + pad)
        y = figh - pad - (r + 1) * ch - r * pad
        _rbox(ax, x, y, cw, ch, fc=PANEL_H, ec=BD_H, r=0.09)
        ax.add_patch(Rectangle((x, y), 0.09, ch, fc=ACC_H, ec="none"))
        icx, icy = x + 0.82, y + ch / 2
        ax.add_patch(Circle((icx, icy), 0.46, fc="white", ec=BD_H, lw=1.0))
        _icon(ax, icon, icx, icy, 0.28)
        ax.text(x + 1.55, y + ch - 0.42, title, fontsize=14, color=INK_H, fontproperties=_BOLD, va="top")
        ax.text(x + 1.55, y + ch - 0.92, textwrap.fill(desc, wrap), fontsize=10.5, color=MUT_H,
                fontproperties=_REG, va="top", linespacing=1.3)
    fig.tight_layout(); fig.savefig(path, dpi=170, bbox_inches="tight", facecolor="white")
    plt.close(fig); return path


def fig_data_sources():
    path = os.path.join(ASSETS, "data_sources.png")
    fig, ax = plt.subplots(figsize=(12, 5.4)); ax.set_xlim(0, 12); ax.set_ylim(0, 5.4); ax.axis("off")
    sources = [(4.05, "Alpaca", "Brokerage & market data", "bars", ALPACA_LOGO,
                "daily prices · option quotes · trading"),
               (2.35, "SEC EDGAR", "Official company filings", "doc", SEC_LOGO,
                "company financials"),
               (0.65, "Ken French Data Library", "Academic factor data", "chartline", None,
                "factor returns for risk")]
    node_x, node_y = 8.7, 2.05
    ends = {4.05: 3.05, 2.35: 2.45, 0.65: 1.85}
    for cy, name, sub, icon, logo, lbl in sources:
        start, end = (4.7, cy + 0.6), (node_x - 0.05, ends[cy])
        ax.add_patch(FancyArrowPatch(start, end, arrowstyle="-|>", mutation_scale=18, color=ACC_H,
                     lw=2.2, shrinkA=0, shrinkB=2, connectionstyle="arc3,rad=0.04"))
        mx, my = (start[0] + end[0]) / 2 + 0.05, (start[1] + end[1]) / 2 + 0.16
        ax.text(mx, my, lbl, fontsize=9.5, color=MUT_H, fontproperties=_REG, ha="center",
                bbox=dict(fc="white", ec="none", pad=1.5))
    for cy, name, sub, icon, logo, lbl in sources:
        _rbox(ax, 0.3, cy, 4.4, 1.2, fc=PANEL_H, ec=BD_H, r=0.1)
        ax.add_patch(Rectangle((0.3, cy), 0.09, 1.2, fc=ACC_H, ec="none"))
        chx, chy = 1.05, cy + 0.6
        _rbox(ax, chx - 0.5, chy - 0.5, 1.0, 1.0, fc="white", ec=BD_H, lw=1.0, r=0.18)
        if logo and os.path.exists(logo):
            _place_logo(ax, logo, chx, chy, 0.4)
        else:
            _icon(ax, icon, chx, chy, 0.32)
        ax.text(1.75, cy + 0.74, name, fontsize=14.5, color=INK_H, fontproperties=_BOLD, va="center")
        ax.text(1.75, cy + 0.36, sub, fontsize=10.5, color=MUT_H, fontproperties=_REG, va="center")
    _rbox(ax, node_x, node_y, 3.0, 1.3, fc=INK_H, ec=INK_H, r=0.12)
    ax.text(node_x + 1.5, node_y + 0.8, "Factor model", fontsize=15, color="white", fontproperties=_BOLD, ha="center")
    ax.text(node_x + 1.5, node_y + 0.4, "& portfolio", fontsize=15, color="white", fontproperties=_BOLD, ha="center")
    ax.add_patch(FancyArrowPatch((node_x + 1.5, node_y - 0.05), (node_x + 1.5, node_y - 0.5),
                 arrowstyle="-|>", mutation_scale=16, color=ACCDK_H, lw=2.2))
    ax.text(node_x + 1.5, node_y - 0.8, "holdings + covered calls", fontsize=10, color=MUT_H,
            fontproperties=_REG, ha="center")
    fig.tight_layout(); fig.savefig(path, dpi=170, bbox_inches="tight", facecolor="white")
    plt.close(fig); return path


def fig_lifecycle():
    path = os.path.join(ASSETS, "lifecycle.png")
    fig, ax = plt.subplots(figsize=(12, 3.2)); ax.set_xlim(0, 12); ax.set_ylim(0, 3.2); ax.axis("off")
    stages = [("Rebalance day", "Sell ~30-day calls\non the new holdings", ACCDK_H),
              ("Through the month", "Hold the stock,\ncollect the premium", INK_H),
              ("Before expiry", "Close the calls\n(early, before earnings)", INK_H),
              ("Next month", "Rewrite fresh calls\non the new holdings", ACCDK_H)]
    n = len(stages); bw = 2.4; gap = (12 - 0.4 - n * bw) / (n - 1); y = 1.35; bh = 1.45
    xs = []
    for i, (t, d, fc) in enumerate(stages):
        x = 0.2 + i * (bw + gap); xs.append(x)
        _rbox(ax, x, y, bw, bh, fc=fc, ec=fc, r=0.12)
        ax.text(x + bw/2, y + bh - 0.34, t, fontsize=12.5, color="white", fontproperties=_BOLD, ha="center")
        ax.text(x + bw/2, y + 0.5, d, fontsize=9.5, color=LIGHT_H, fontproperties=_REG, ha="center", linespacing=1.25)
        if i < n - 1:
            ax.add_patch(FancyArrowPatch((x + bw + 0.06, y + bh/2), (x + bw + gap - 0.06, y + bh/2),
                         arrowstyle="-|>", mutation_scale=16, color=ACC_H, lw=2.4))
    ax.add_patch(FancyArrowPatch((xs[-1] + bw/2, y - 0.04), (xs[0] + bw/2, y - 0.04),
                 connectionstyle="arc3,rad=0.34", arrowstyle="-|>", mutation_scale=15, color=MUT_H, lw=1.8))
    ax.text(6, 0.22, "repeats every month", fontsize=11, color=MUT_H, fontproperties=_BOLD, ha="center")
    fig.tight_layout(); fig.savefig(path, dpi=170, bbox_inches="tight", facecolor="white")
    plt.close(fig); return path


def fig_methodology():
    path = os.path.join(ASSETS, "methodology.png")
    items = [("calendar", "Walk-forward", "Tested month by month, using only data known then"),
             ("percent", "Realistic costs", "A trading cost charged on every simulated trade"),
             ("chartline", "Real option prices", "Actual historical premiums used where available"),
             ("check", "Index cross-check", "Reproduces the public BuyWrite index closely")]
    fig, ax = plt.subplots(figsize=(12, 3.0)); ax.set_xlim(0, 12); ax.set_ylim(0, 3.0); ax.axis("off")
    n = 4; pad = 0.3; cw = (12 - pad * (n + 1)) / n; y = 0.2; ch = 2.6
    for i, (icon, t, d) in enumerate(items):
        x = pad + i * (cw + pad)
        _rbox(ax, x, y, cw, ch, fc=PANEL_H, ec=BD_H, r=0.09)
        ax.add_patch(Rectangle((x, y + ch - 0.09), cw, 0.09, fc=ACC_H, ec="none"))
        icx, icy = x + cw/2, y + ch - 0.72
        ax.add_patch(Circle((icx, icy), 0.46, fc="white", ec=BD_H, lw=1.0)); _icon(ax, icon, icx, icy, 0.28)
        ax.text(x + cw/2, y + ch - 1.5, t, fontsize=13, color=INK_H, fontproperties=_BOLD, ha="center")
        ax.text(x + cw/2, y + ch - 1.82, textwrap.fill(d, 26), fontsize=10, color=MUT_H,
                fontproperties=_REG, ha="center", va="top", linespacing=1.3)
    fig.tight_layout(); fig.savefig(path, dpi=170, bbox_inches="tight", facecolor="white")
    plt.close(fig); return path


def fig_risk_cards():
    return _icon_cards(os.path.join(ASSETS, "risk_cards.png"), [
        ("shield", "Pre-trade check", "No order can break a limit, and no call can be sold uncovered"),
        ("refresh", "Safe to re-run", "Can stop and restart mid-rebalance without double-trading"),
        ("stop", "Emergency stop", "One command halts trading or sells the whole book to cash"),
        ("doc", "Full audit log", "Every score, order, and fill is recorded"),
    ], cols=2)


def fig_integrity_cards():
    return _icon_cards(os.path.join(ASSETS, "integrity_cards.png"), [
        ("split", "Corporate actions", "Splits and dividends never look like returns"),
        ("clock", "Point-in-time", "Only data public on each date is used, with no hindsight"),
        ("refresh", "Reconciled", "Squared against the broker on every run"),
        ("funnel", "Clean universe", "Only liquid stocks with reliable filings are eligible"),
    ], cols=2)


def fig_dashboard():
    """A schematic of the live dashboard, in its actual dark theme (illustrative, no live figures)."""
    path = os.path.join(ASSETS, "dashboard.png")
    BG, PANEL, BORD = "#0B1322", "#172439", "#28384f"
    TEAL, MUT, TXT, FAINT = "#46B8AD", "#8A97AC", "#E6ECF2", "#34425e"
    fig, ax = plt.subplots(figsize=(12, 4.9)); ax.set_xlim(0, 12); ax.set_ylim(0, 4.9); ax.axis("off")
    _rbox(ax, 0.2, 0.2, 11.6, 4.5, fc=BG, ec=BORD, lw=1.4, r=0.04)
    gx, gy = 0.62, 4.32
    for rr in range(3):
        for cc in range(3):
            ax.add_patch(Circle((gx + cc*0.13, gy - rr*0.13), 0.045, fc=(TEAL if cc == 1 else FAINT), ec="none"))
    ax.text(1.3, 4.18, "Systematic Factor Income Fund", fontsize=12, color=TXT, fontproperties=_BOLD, va="center")
    ax.add_patch(Circle((9.45, 4.18), 0.06, fc=TEAL, ec="none"))
    ax.text(9.6, 4.18, "streaming · live", fontsize=9.5, color=TEAL, fontproperties=_REG, va="center")
    ax.plot([0.2, 11.8], [3.92, 3.92], color=BORD, lw=1)
    tx = 0.62
    for name, active in [("Overview", True), ("Portfolio", False), ("Performance", False), ("Backtest", False)]:
        ax.text(tx, 3.62, name, fontsize=10.5, color=(TEAL if active else MUT),
                fontproperties=(_BOLD if active else _REG), va="center")
        w = len(name) * 0.118 + 0.05
        if active:
            ax.plot([tx, tx + w], [3.4, 3.4], color=TEAL, lw=2.5)
        tx += w + 0.55
    for i, lab in enumerate(["NAV", "Day P&L", "Leverage", "Premium"]):
        x = 0.62 + i * 2.75
        _rbox(ax, x, 2.2, 2.55, 0.95, fc=PANEL, ec=BORD, lw=1, r=0.08)
        ax.text(x + 0.22, 2.87, lab, fontsize=9, color=MUT, fontproperties=_REG, va="center")
        _rbox(ax, x + 0.22, 2.42, 1.5, 0.2, fc=(TEAL if i == 0 else FAINT), ec=(TEAL if i == 0 else FAINT), r=0.1)
    _rbox(ax, 0.62, 0.45, 7.2, 1.5, fc=PANEL, ec=BORD, lw=1, r=0.05)
    xs = np.linspace(0.95, 7.5, 60); base = np.linspace(0, 1, 60)
    ax.plot(xs, 0.72 + (base + 0.1 * np.sin(np.linspace(0, 7, 60))) * 0.95, color=TEAL, lw=2)
    ax.plot(xs, 0.72 + base * 0.6, color=MUT, lw=1.4, ls=(0, (4, 3)))
    ax.text(0.85, 1.78, "Growth vs benchmark", fontsize=8.5, color=MUT, fontproperties=_REG)
    cx, cy, start = 10.0, 1.15, 90
    for fr, co in [(0.32, TEAL), (0.27, "#2E8B82"), (0.23, "#3A4C63"), (0.18, "#26354e")]:
        ax.add_patch(Wedge((cx, cy), 0.6, start, start + fr * 360, width=0.24, fc=co, ec="none"))
        start += fr * 360
    ax.text(cx, cy - 0.95, "Allocation", fontsize=8.5, color=MUT, ha="center", fontproperties=_REG)
    fig.tight_layout(); fig.savefig(path, dpi=170, bbox_inches="tight", transparent=True)
    plt.close(fig); return path


print("rendering assets…")
F_Z = formula(r"z=\dfrac{x-\mu}{\sigma}", "f_z.png", 44)
F_VALUE = formula(r"\mathrm{Value}=z\left(\,z(E/P)+z(B/P)\,\right)", "f_value.png", 40)
F_VALUE_W = formula(r"E/P=\dfrac{1}{P/E}\qquad B/P=\dfrac{1}{P/B}", "f_value_w.png", 30, MUT_H)
F_MOM = formula(r"\mathrm{Momentum}=z\left(\dfrac{P_{t-21}}{P_{t-252}}-1\right)", "f_mom.png", 40)
F_QUAL = formula(r"\mathrm{Quality}=z\left(\,z(\mathrm{ROE})+z(\mathrm{margin})\,\right)", "f_qual.png", 40)
F_LOWVOL = formula(r"\mathrm{Low\ volatility}=z\left(-\,\sigma\right)", "f_lowvol.png", 40)
F_SCORE = formula(r"\mathrm{Score}=\dfrac{1}{4}\left(Q+V+M+L\right)", "f_score.png", 42)
F_OBJ = formula(r"\max\ \sum_i w_i\cdot\mathrm{score}_i", "f_obj.png", 42)
IMG_BELL = fig_bell()
IMG_PAYOFF = fig_payoff()
IMG_RR = fig_risk_return()
IMG_VRP = fig_vrp()
IMG_SOURCES = fig_data_sources()
IMG_LIFECYCLE = fig_lifecycle()
IMG_METHOD = fig_methodology()
IMG_RISK = fig_risk_cards()
IMG_INTEGRITY = fig_integrity_cards()
IMG_DASH = fig_dashboard()


# =========================================================================== #
# Deck scaffold + helpers
# =========================================================================== #
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
EMU_W, EMU_H = prs.slide_width, prs.slide_height
_n = {"v": 0}


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(s, l, t, w, h, runs, size=16, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.1, font=FONT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, bold, color)]
    for r in runs:
        if isinstance(r, str):
            rt, rb, rc = r, bold, color
        else:
            rt = r[0]; rb = r[1] if len(r) > 1 else bold; rc = r[2] if len(r) > 2 else color
        run = p.add_run(); run.text = rt
        run.font.name = font; run.font.size = Pt(size); run.font.bold = rb; run.font.color.rgb = rc
    return tb


def bullets(s, items, l, t, w, h, size=16, gap=11, line_spacing=1.12):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    first = True
    for it in items:
        lvl = it.get("lvl", 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing; p.space_after = Pt(it.get("gap", gap))
        g = p.add_run(); g.text = "▪  " if lvl == 0 else "        ·  "
        g.font.name = FONT; g.font.size = Pt(size if lvl == 0 else size - 2)
        g.font.bold = lvl == 0; g.font.color.rgb = ACCENT if lvl == 0 else MUTED
        base = it.get("color", TEXT if lvl == 0 else MUTED)
        for r in it["runs"]:
            if isinstance(r, str):
                rt, rb, rc = r, False, base
            else:
                rt = r[0]; rb = r[1] if len(r) > 1 else False; rc = r[2] if len(r) > 2 else base
            run = p.add_run(); run.text = rt
            run.font.name = FONT; run.font.size = Pt(size if lvl == 0 else size - 2)
            run.font.bold = rb; run.font.color.rgb = rc
    return tb


def arrow(s, l, t, w, h, color=ACCENT):
    sp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    sp.adjustments[0] = 0.55; sp.adjustments[1] = 0.6
    return sp


def image_fit(s, path, l, t, w, h, halign="center", valign="middle"):
    pic = s.shapes.add_picture(path, l, t)
    scale = min(w / pic.width, h / pic.height)
    pic.width = int(pic.width * scale); pic.height = int(pic.height * scale)
    pic.left = int(l + (w - pic.width) / 2) if halign == "center" else (
        int(l + w - pic.width) if halign == "right" else l)
    pic.top = int(t + (h - pic.height) / 2) if valign == "middle" else (
        int(t + h - pic.height) if valign == "bottom" else t)
    return pic


def slide():
    return prs.slides.add_slide(BLANK)


def header(s, kicker, title, subtitle=None):
    rect(s, 0, 0, EMU_W, EMU_H, fill=WHITE)
    rect(s, Inches(0.7), Inches(0.62), Inches(0.34), Inches(0.11), fill=ACCENT)
    text(s, Inches(1.12), Inches(0.55), Inches(11), Inches(0.3), kicker.upper(), size=12,
         color=ACCENT_DK, bold=True)
    text(s, Inches(0.7), Inches(0.92), Inches(12), Inches(0.7), title, size=26, color=INK, bold=True)
    y = 1.6
    if subtitle:
        text(s, Inches(0.7), Inches(1.52), Inches(12.3), Inches(0.5), subtitle, size=15, color=MUTED)
        y = 2.02
    rect(s, Inches(0.7), Inches(y), Inches(11.93), Pt(1.2), fill=RULE)
    return y + 0.2


def footer(s):
    _n["v"] += 1
    text(s, Inches(0.7), Inches(7.04), Inches(8), Inches(0.3),
         "Systematic Factor Income Fund    ·    confidential", size=9, color=MUTED)
    text(s, Inches(11.0), Inches(7.04), Inches(1.63), Inches(0.3), f"{_n['v']:02d}", size=9,
         color=MUTED, align=PP_ALIGN.RIGHT)


def glyph(s, left, top, d=0.17, gap=0.29, faint=RGBColor(0x3A, 0x4C, 0x63)):
    """The SFI brand mark: a 3x3 grid of dots, centre column teal (for dark backgrounds)."""
    for rr in range(3):
        for cc in range(3):
            rect(s, Inches(left + cc * gap), Inches(top + rr * gap), Inches(d), Inches(d),
                 fill=(ACCENT if cc == 1 else faint), shape=MSO_SHAPE.OVAL)


def divider(num, title, subtitle=None):
    s = slide(); rect(s, 0, 0, EMU_W, EMU_H, fill=INK)
    glyph(s, 12.0, 0.62, d=0.15, gap=0.26)
    text(s, Inches(0.6), Inches(1.0), Inches(7), Inches(3.2), f"{num:02d}", size=200,
         color=SLATE, bold=True)
    rect(s, Inches(0.95), Inches(4.35), Inches(2.4), Pt(3), fill=ACCENT)
    text(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(1.0), title, size=40, color=WHITE, bold=True)
    if subtitle:
        text(s, Inches(0.95), Inches(5.7), Inches(11), Inches(1.0), subtitle, size=16, color=LIGHT)
    return s


def chip(s, x, y, w, h, head, sub, head_color=INK):
    rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=PANEL, line=PANEL_BD, line_w=1,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    rect(s, Inches(x), Inches(y), Inches(0.09), Inches(h), fill=ACCENT)
    text(s, Inches(x + 0.24), Inches(y + 0.12), Inches(w - 0.34), Inches(0.3), head, size=13.5,
         color=head_color, bold=True)
    text(s, Inches(x + 0.24), Inches(y + 0.46), Inches(w - 0.34), Inches(h - 0.5), sub, size=10.5,
         color=MUTED, line_spacing=1.05)


def formula_card(s, x, y, w, h, img, label=None, where_img=None, img_h_frac=0.62):
    rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=PANEL, line=PANEL_BD, line_w=1,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    top = y
    if label:
        text(s, Inches(x + 0.3), Inches(y + 0.16), Inches(w - 0.6), Inches(0.3), label.upper(),
             size=11, color=ACCENT_DK, bold=True)
        top = y + 0.5
    fh = (y + h) - top - (0.0 if not where_img else 0.55)
    image_fit(s, img, Inches(x + 0.3), Inches(top), Inches(w - 0.6), Inches(fh), valign="middle")
    if where_img:
        image_fit(s, where_img, Inches(x + 0.3), Inches(y + h - 0.62), Inches(w - 0.6),
                  Inches(0.5), valign="middle")


def content(kicker, title, items, subtitle=None, size=16, body_top=None, body_w=11.93):
    s = slide(); y = header(s, kicker, title, subtitle)
    bullets(s, items, Inches(0.7), Inches(body_top or y), Inches(body_w), Inches(5.0), size=size)
    footer(s); return s


def two_col(kicker, title, lhead, litems, rhead, ritems, subtitle=None):
    s = slide(); y = header(s, kicker, title, subtitle)
    text(s, Inches(0.7), Inches(y + 0.05), Inches(5.7), Inches(0.4), lhead, size=15, color=INK, bold=True)
    bullets(s, litems, Inches(0.7), Inches(y + 0.6), Inches(5.7), Inches(4.4), size=14)
    rect(s, Inches(6.66), Inches(y + 0.1), Pt(1.2), Inches(4.6), fill=RULE)
    text(s, Inches(6.95), Inches(y + 0.05), Inches(5.7), Inches(0.4), rhead, size=15, color=INK, bold=True)
    bullets(s, ritems, Inches(6.95), Inches(y + 0.6), Inches(5.7), Inches(4.4), size=14)
    footer(s); return s


def table_slide(kicker, title, headers, rows, col_w, subtitle=None, note=None, hl=None,
                fsize=12, hsize=12):
    s = slide(); y = header(s, kicker, title, subtitle)
    hl = hl or set()
    nrows = len(rows) + 1
    gfx = s.shapes.add_table(nrows, len(headers), Inches(0.7), Inches(y + 0.15),
                             Inches(sum(col_w)), Inches(0.47 * nrows))
    tbl = gfx.table; tbl.first_row = False; tbl.horz_banding = False
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(cw)
    for c, head in enumerate(headers):
        cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = INK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(cell, m, Inches(0.07))
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = head; r.font.name = FONT; r.font.size = Pt(hsize)
        r.font.bold = True; r.font.color.rgb = WHITE
    for ri, row in enumerate(rows, 1):
        hot = (ri - 1) in hl
        for c, val in enumerate(row):
            cell = tbl.cell(ri, c); cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xEC, 0xF6, 0xF4) if hot else (PANEL if ri % 2 == 0 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                setattr(cell, m, Inches(0.07))
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val); r.font.name = FONT; r.font.size = Pt(fsize)
            r.font.bold = hot or c == 0; r.font.color.rgb = INK if (hot or c == 0) else TEXT
    if note:
        text(s, Inches(0.7), Inches(y + 0.15 + 0.47 * nrows + 0.15), Inches(11.93), Inches(0.9),
             note, size=11, color=MUTED, line_spacing=1.15)
    footer(s); return s


def stat_slide(kicker, title, stats, subtitle=None, caption=None):
    s = slide(); y = header(s, kicker, title, subtitle)
    n = len(stats); gap = 0.3; cw = (11.93 - gap * (n - 1)) / n; top = y + 0.6
    for i, st in enumerate(stats):
        x = 0.7 + i * (cw + gap)
        rect(s, Inches(x), Inches(top), Inches(cw), Inches(2.0), fill=PANEL, line=PANEL_BD,
             line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        rect(s, Inches(x), Inches(top), Inches(cw), Inches(0.09), fill=ACCENT)
        text(s, Inches(x), Inches(top + 0.42), Inches(cw), Inches(0.9), st[0], size=38,
             color=st[2] if len(st) > 2 else INK, bold=True, align=PP_ALIGN.CENTER)
        text(s, Inches(x + 0.15), Inches(top + 1.35), Inches(cw - 0.3), Inches(0.6), st[1],
             size=12, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.05)
    if caption:
        text(s, Inches(0.7), Inches(top + 2.35), Inches(11.93), Inches(1.0), caption, size=12.5,
             color=MUTED, line_spacing=1.2)
    footer(s); return s


# =========================================================================== #
# BUILD
# =========================================================================== #
# ---- Cover ---------------------------------------------------------------- #
s = slide(); rect(s, 0, 0, EMU_W, EMU_H, fill=INK); rect(s, 0, 0, Inches(0.22), EMU_H, fill=ACCENT)
glyph(s, 11.5, 0.92)
text(s, Inches(0.9), Inches(0.95), Inches(10), Inches(0.4), "PROPRIETARY INVESTMENT STRATEGY",
     size=13, color=ACCENT, bold=True)
text(s, Inches(0.85), Inches(2.3), Inches(11.9), Inches(2.0), "Systematic Factor\nIncome Fund",
     size=52, color=WHITE, bold=True, line_spacing=1.02)
rect(s, Inches(0.9), Inches(4.55), Inches(2.6), Pt(2.5), fill=ACCENT)
text(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(0.6),
     "A rules-based equity portfolio with an options income overlay", size=20, color=LIGHT)
text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
     [("Strategy overview", True, WHITE),
      ("        Confidential, for internal review", False, RGBColor(0x8B, 0x99, 0xA8))], size=12)

# ---- Agenda --------------------------------------------------------------- #
s = slide(); y = header(s, "Contents", "What this overview covers")
agenda = [
    ("1", "The objective", "What the fund is built to do"),
    ("2", "Why this approach", "Proven return drivers, plus an income layer"),
    ("3", "The factor strategy", "The four traits we select stocks on"),
    ("4", "Building the portfolio", "Turning scores into holdings, within strict limits"),
    ("5", "The income overlay", "Earning option premium on what we own"),
    ("6", "The data pipeline", "Where the numbers come from"),
    ("7", "Execution, risk and results", "How we trade, how we control risk, how it tested"),
    ("8", "Operations", "Monitoring and infrastructure"),
    ("9", "Status and outlook", "Where the build stands today"),
    ("·", "Appendix", "Formulas, sources, detailed results, glossary"),
]
top = y + 0.18
for i, (num, t_, d_) in enumerate(agenda):
    ry = top + i * 0.485
    text(s, Inches(0.7), Inches(ry), Inches(0.6), Inches(0.4), num, size=15, color=ACCENT_DK, bold=True)
    text(s, Inches(1.4), Inches(ry), Inches(3.5), Inches(0.4), t_, size=15, color=INK, bold=True)
    text(s, Inches(5.0), Inches(ry), Inches(7.6), Inches(0.4), d_, size=13, color=MUTED)
footer(s)

# =========================================================================== #
# § 1 — THE OBJECTIVE
# =========================================================================== #
divider(1, "The objective", "A book built to diversify the firm with steady, risk-adjusted returns")

content("1 · The objective", "What the fund is built to do", [
    {"runs": [("Preserve capital first. ", True), ("Every design choice favours protecting the book "
      "over maximising return. Shallow drawdowns matter more than big years.")]},
    {"runs": [("Produce uncorrelated returns. ", True), ("A return stream that diversifies the firm's "
      "core business rather than echoing it.")]},
    {"runs": [("Stay explainable. ", True), ("A transparent, rules-based strategy with a long research "
      "record. Anyone can audit why a position is held.")]},
    {"runs": [("Earn income steadily. ", True), ("Option premium provides a recurring return that does "
      "not depend on the market rising.")]},
], subtitle="The mandate is risk-adjusted income and capital preservation, not maximum return")

s = slide(); y = header(s, "1 · The objective", "What it is, and what it is not",
                        "Stating the boundaries up front keeps the strategy disciplined")
left = [("A systematic factor portfolio", "Rules select stocks on proven traits, rebalanced monthly."),
        ("An options income layer", "Covered calls turn holdings into a recurring income stream."),
        ("Broadly diversified", "About 20 stocks, with firm limits on any one name or sector.")]
right = [("Not a forecasting model", "No black-box prediction of next month's returns."),
         ("Not market timing", "No attempt to call tops or bottoms. The factors carry the defense."),
         ("Not concentrated bets", "No large single-stock positions or thematic gambles.")]
text(s, Inches(0.7), Inches(y + 0.05), Inches(5.8), Inches(0.4), "WHAT IT IS", size=13,
     color=ACCENT_DK, bold=True)
for i, (h_, d_) in enumerate(left):
    chip(s, 0.7, y + 0.5 + i * 1.35, 5.7, 1.2, h_, d_)
text(s, Inches(6.95), Inches(y + 0.05), Inches(5.8), Inches(0.4), "WHAT IT IS NOT", size=13,
     color=MUTED, bold=True)
for i, (h_, d_) in enumerate(right):
    chip(s, 6.95, y + 0.5 + i * 1.35, 5.7, 1.2, h_, d_, head_color=INK)
footer(s)

# Two engines
s = slide(); y = header(s, "1 · The objective", "Two engines, one portfolio",
                        "A stock portfolio decides what we own. An overlay earns income on it.")
rect(s, Inches(0.9), Inches(2.7), Inches(4.7), Inches(1.95), fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
text(s, Inches(0.95), Inches(3.05), Inches(4.6), Inches(0.5), "FACTOR STOCK PORTFOLIO", size=15,
     color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1.1), Inches(3.7), Inches(4.3), Inches(0.8),
     "About 20 stocks chosen on four\nproven traits, rebalanced monthly", size=12, color=LIGHT,
     align=PP_ALIGN.CENTER, line_spacing=1.1)
rect(s, Inches(7.75), Inches(2.7), Inches(4.7), Inches(1.95), fill=ACCENT_DK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
text(s, Inches(7.8), Inches(3.05), Inches(4.6), Inches(0.5), "COVERED-CALL OVERLAY", size=15,
     color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(7.95), Inches(3.7), Inches(4.3), Inches(0.8),
     "Sell call options on those holdings\nto collect cash premium", size=12, color=RGBColor(0xCC, 0xEE, 0xE9),
     align=PP_ALIGN.CENTER, line_spacing=1.1)
arrow(s, Inches(5.85), Inches(3.45), Inches(1.7), Inches(0.45))
text(s, Inches(5.7), Inches(3.0), Inches(2.0), Inches(0.4), "earns income on", size=11, color=MUTED,
     align=PP_ALIGN.CENTER)
text(s, Inches(0.9), Inches(5.2), Inches(11.6), Inches(1.2),
     [("Total return comes from two sources: ", True),
      ("the stocks we select, and the premium the overlay collects. The overlay also lowers the "
       "portfolio's volatility, in exchange for capping gains on any stock that rises sharply.")],
     size=15, line_spacing=1.3)
footer(s)

# =========================================================================== #
# § 2 — WHY THIS APPROACH
# =========================================================================== #
divider(2, "Why this approach", "Built on durable, well-documented sources of return")

content("2 · Why this approach", "Why proven factors, not prediction", [
    {"runs": [("Predicting short-term returns is the most crowded game in investing. ", True),
      ("Thousands of funds chase the same price signals, and any edge erodes quickly.")]},
    {"runs": [("Factors are paid for bearing risk, not for finding mistakes. ", True),
      ("That is why they have persisted across decades and across markets, rather than being "
       "arbitraged away.")]},
    {"runs": [("A rules-based portfolio is defensible. ", True),
      ("It can be explained and audited, and it holds up to scrutiny during a weak stretch. A "
       "black-box model cannot.")]},
    {"runs": [("It fits the mandate. ", True),
      ("The goal is durable, risk-adjusted income, which is exactly what factor premia and an "
       "options overlay are designed to provide.")]},
], subtitle="The strategy relies on effects with a long academic and practical record")

s = slide(); y = header(s, "2 · Why this approach", "Why an options income overlay",
                        "Options consistently trade richer than the price moves that follow")
bullets(s, [
    {"runs": [("Implied volatility is the price movement the market expects, ", True),
      ("and it is built into every option's price.")]},
    {"runs": [("It usually runs higher than the movement that actually occurs. ", True),
      ("Sellers of options are paid for that gap. It is a persistent, well-studied premium.")]},
    {"runs": [("We collect it on stock we already own. ", True),
      ("Selling calls on our holdings turns that premium into income, with no extra capital tied up.")]},
    {"runs": [("It also smooths the ride. ", True),
      ("Premium cushions down months and lowers overall volatility, at the cost of capping the "
       "occasional very large gain.")]},
], Inches(0.7), Inches(y), Inches(6.1), Inches(4.4), size=14.5)
image_fit(s, IMG_VRP, Inches(7.1), Inches(y + 0.1), Inches(5.5), Inches(4.3), valign="top")
footer(s)

# =========================================================================== #
# § 3 — THE FACTOR STRATEGY
# =========================================================================== #
divider(3, "The factor strategy", "Four traits, measured the same way, combined into one score")

# What is a factor (diagram)
s = slide(); y = header(s, "3 · The factor strategy", "What we mean by a factor",
                        "A factor is a company trait that has historically earned higher long-run returns")
text(s, Inches(0.7), Inches(y), Inches(11.93), Inches(0.85),
     "We score every stock on four such traits, combine the four scores into one, and hold the "
     "highest-ranked names. Each trait below is supported by decades of research.", size=15,
     color=TEXT, line_spacing=1.25)
dy = y + 1.1
facs = [("Quality", "profitable, efficient companies"), ("Value", "inexpensive vs. fundamentals"),
        ("Momentum", "sustained recent strength"), ("Low volatility", "steadier price behaviour")]
for i, (h_, sub) in enumerate(facs):
    chip(s, 0.7, dy + i * 0.92, 4.0, 0.78, h_, sub)
arrow(s, Inches(4.95), Inches(dy + 1.55), Inches(0.7), Inches(0.5))
rect(s, Inches(5.85), Inches(dy + 0.9), Inches(3.1), Inches(1.85), fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
text(s, Inches(5.95), Inches(dy + 1.2), Inches(2.9), Inches(1.3), "Combined\nscore", size=18,
     color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
arrow(s, Inches(9.15), Inches(dy + 1.55), Inches(0.7), Inches(0.5))
rect(s, Inches(10.0), Inches(dy + 0.9), Inches(2.63), Inches(1.85), fill=ACCENT_DK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
text(s, Inches(10.1), Inches(dy + 1.2), Inches(2.43), Inches(1.3), "Hold the\nbest ~20", size=18,
     color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer(s)

# z-score
s = slide(); y = header(s, "3 · The factor strategy", "Putting every company on the same scale",
                        "Each trait becomes a score showing how far a company sits from the average")
text(s, Inches(0.7), Inches(y), Inches(6.0), Inches(2.1),
     [("A score (a z-score) measures how far a company sits above or below the average, counted in "
       "standard deviations. A standard deviation is the usual gauge of how spread out a set of "
       "numbers is.\n\n", False, TEXT),
      ("A score of +1 means one standard deviation better than average. Zero is exactly average.",
       False, TEXT)], size=15, line_spacing=1.25)
formula_card(s, 0.7, y + 2.4, 6.0, 1.35, F_Z, label="How a score is calculated")
text(s, Inches(0.95), Inches(y + 3.95), Inches(5.8), Inches(0.5),
     [("x", True, INK), (" the company's value    ", False, MUTED),
      ("μ", True, INK), (" the average    ", False, MUTED),
      ("σ", True, INK), (" standard deviation", False, MUTED)], size=12.5)
image_fit(s, IMG_BELL, Inches(7.0), Inches(y + 0.1), Inches(5.6), Inches(4.5), valign="top")
footer(s)


def factor_slide(kicker_title, intuition, defs, img, where_img=None, source=None):
    s = slide(); y = header(s, "3 · The four factors", kicker_title)
    text(s, Inches(0.7), Inches(y), Inches(11.93), Inches(0.6), intuition, size=15.5, color=TEXT,
         line_spacing=1.2)
    bullets(s, defs, Inches(0.7), Inches(y + 0.85), Inches(11.93), Inches(2.0), size=14)
    cardh = 1.55 if where_img else 1.4
    cy = 5.0 if where_img else 5.15
    formula_card(s, 0.7, cy, 11.93, cardh, img, label="How it is measured", where_img=where_img)
    if source:
        text(s, Inches(0.7), Inches(cy + cardh + 0.12), Inches(11.93), Inches(0.4),
             [("Foundational research:  ", True, MUTED), (source, False, MUTED)], size=10.5)
    footer(s); return s


factor_slide(
    "Value: paying less for what you get",
    "Cheaper stocks, judged against company fundamentals, have historically out-returned expensive ones.",
    [{"runs": [("Price-to-earnings (P/E): ", True), ("share price divided by yearly earnings per share. "
       "It shows how many dollars you pay for one dollar of profit. Lower is cheaper.")]},
     {"runs": [("Price-to-book (P/B): ", True), ("share price divided by net asset value per share. "
       "Lower is cheaper.")]},
     {"runs": [("We score the inverse (the yield), ", True), ("so a lower price produces a higher Value score.")]}],
    F_VALUE, where_img=F_VALUE_W,
    source="Fama & French (1992, 1993); Basu (1977).")

factor_slide(
    "Momentum: recent strength tends to continue",
    "Stocks that have outperformed over the past year have tended to keep outperforming over the next months.",
    [{"runs": [("We measure the price change over the past year, ", True), ("skipping the most recent "
       "month, which tends to reverse.")]},
     {"runs": [("P is the share price; ", True), ("the subscripts count trading days back from today "
       "(252 days is about one year, 21 is about one month).")]},
     {"runs": [("The result is then put on the same 0-centred scale ", True), ("as the other factors.")]}],
    F_MOM,
    source="Jegadeesh & Titman (1993); Asness, Moskowitz & Pedersen (2013).")

factor_slide(
    "Quality: profitable, well-run businesses",
    "Highly profitable, efficient companies have delivered better risk-adjusted returns over time.",
    [{"runs": [("Return on equity (ROE): ", True), ("annual profit as a percentage of shareholders' "
       "money. Higher is better.")]},
     {"runs": [("Margin (gross margin): ", True), ("the share of revenue left after the direct cost of "
       "goods. Higher means more pricing power.")]},
     {"runs": [("Quality companies also tend to hold up better in downturns, ", True),
       ("which adds a defensive tilt without any market timing.")]}],
    F_QUAL,
    source="Novy-Marx (2013); Asness, Frazzini & Pedersen (2019); Piotroski (2000).")

factor_slide(
    "Low volatility: steadier stocks, better risk-adjusted returns",
    "Calmer stocks have historically delivered stronger returns per unit of risk than the most volatile ones.",
    [{"runs": [("Volatility (σ): ", True), ("how much a stock's daily price bounces around, measured "
       "over the past year. Higher means a wilder ride.")]},
     {"runs": [("We score the negative of volatility, ", True), ("so a calmer stock receives a higher score.")]},
     {"runs": [("Together with Quality, ", True), ("this is what makes the portfolio defensive on its "
       "own, without a separate risk-off switch.")]}],
    F_LOWVOL,
    source="Ang, Hodrick, Xing & Zhang (2006); Frazzini & Pedersen (2014); Baker, Bradley & Wurgler (2011).")

# Composite + discipline
s = slide(); y = header(s, "3 · The factor strategy", "Combining the four into one score",
                        "Equal weight, with safeguards that keep the comparison fair and honest")
bullets(s, [
    {"runs": [("Equal weight to start. ", True), ("Each factor counts for one quarter of the score. "
      "With no strong reason to favour one, equal weighting is the most defensible starting point.")]},
    {"runs": [("Extreme values are trimmed before scoring. ", True), ("A handful of outliers cannot "
      "distort the scale, so all four factors genuinely count equally.")]},
    {"runs": [("Only public information is used. ", True), ("A company's financials enter the score "
      "only once they have actually been filed, so the historical test never sees the future.")]},
], Inches(0.7), Inches(y), Inches(6.1), Inches(3.5), size=14.5)
formula_card(s, 7.1, y + 0.7, 5.5, 1.5, F_SCORE, label="The combined score")
text(s, Inches(7.1), Inches(y + 2.45), Inches(5.5), Inches(0.6),
     "Q, V, M, L are the four factor scores. The highest combined scores become the portfolio.",
     size=12, color=MUTED, line_spacing=1.15)
footer(s)

# =========================================================================== #
# § 4 — BUILDING THE PORTFOLIO
# =========================================================================== #
divider(4, "Building the portfolio", "Turning scores into holdings, inside strict risk limits")

s = slide(); y = header(s, "4 · Building the portfolio", "From scores to holdings",
                        "The scores choose the names. Hard limits keep the portfolio diversified.")
bullets(s, [
    {"runs": [("Rank by combined score. ", True), ("The top-scoring names are the candidates to hold.")]},
    {"runs": [("Maximise total score, subject to limits. ", True), ("A standard optimiser sets the "
      "weights to favour the best-scoring names while respecting every limit below.")]},
    {"runs": [("The result is roughly equal-weight across about 20 stocks. ", True),
      ("Equal-weight factor portfolios have travelled well in out-of-sample testing.")]},
], Inches(0.7), Inches(y), Inches(6.1), Inches(3.2), size=14.5)
formula_card(s, 7.1, y + 0.65, 5.5, 1.35, F_OBJ, label="What the optimiser solves")
text(s, Inches(7.1), Inches(y + 2.25), Inches(5.5), Inches(0.7),
     "w is each stock's weight; score is its combined factor score. Read: choose weights that put "
     "the most money in the best-scoring names.", size=12, color=MUTED, line_spacing=1.15)
footer(s)

s = slide(); y = header(s, "4 · Building the portfolio", "The limits are the risk controls",
                        "Diversification is enforced by rule, not left to chance")
limits = [("5% maximum per stock", "No single name can dominate. This is what sets the holding count near 20."),
          ("30% maximum per sector", "No single industry can take over the book."),
          ("~20 holdings, near equal weight", "Broad enough to diversify, focused enough to matter."),
          ("95% invested, 5% cash", "A stable allocation at all times, with a small cash buffer.")]
for i, (h_, d_) in enumerate(limits):
    col = i % 2; row = i // 2
    chip(s, 0.7 + col * 6.13, y + 0.2 + row * 1.55, 5.7, 1.35, h_, d_)
footer(s)

# =========================================================================== #
# § 5 — THE INCOME OVERLAY
# =========================================================================== #
divider(5, "The income overlay", "Earning option premium on the stocks we already own")

s = slide(); y = header(s, "5 · The income overlay", "Selling covered calls",
                        "We sell call options on our holdings and keep the cash premium")
bullets(s, [
    {"runs": [("Call option: ", True), ("a contract that lets its buyer purchase our stock at a fixed "
      "“strike” price before a set date. We sell that right.")]},
    {"runs": [("Premium: ", True), ("the cash the buyer pays us up front. We keep it in every outcome.")]},
    {"runs": [("If the stock stays below the strike, ", True), ("the option expires worthless and we "
      "simply keep the premium and the stock.")]},
    {"runs": [("If it rises above the strike, ", True), ("the stock is sold at the strike. We keep the "
      "premium and the gain up to that point, but give up the gain beyond it.")]},
], Inches(0.7), Inches(y), Inches(6.0), Inches(4.3), size=14)
image_fit(s, IMG_PAYOFF, Inches(6.95), Inches(y + 0.1), Inches(5.75), Inches(4.2), valign="top")
footer(s)

content("5 · The income overlay", "Choosing which call to sell", [
    {"runs": [("Delta, about 0.30. ", True), ("Delta estimates the chance the option ends up being "
      "exercised. Targeting 0.30 means roughly a 30% chance the stock is sold at the strike. Low "
      "enough to keep most of the upside, high enough to collect meaningful premium.")]},
    {"runs": [("Strikes set by probability, not a fixed percentage. ", True), ("Using delta adapts the "
      "strike to each stock automatically. A calm stock and a volatile one get the same ~30% odds.")]},
    {"runs": [("Expiring in 30 to 45 days. ", True), ("Long enough to collect a worthwhile premium, "
      "short enough to refresh monthly with the rest of the portfolio.")]},
    {"runs": [("Sold at the quoted middle price. ", True), ("We take a fair price from the live option "
      "market rather than chasing.")]},
], subtitle="The rule is consistent across every stock, with no discretion")

s = slide(); y = header(s, "5 · The income overlay", "Managing the options through the month",
                        "One simple monthly cycle, repeated against each new set of holdings")
image_fit(s, IMG_LIFECYCLE, Inches(0.7), Inches(y + 0.25), Inches(11.93), Inches(3.2), valign="top")
bullets(s, [
    {"runs": [("Closed early before earnings, then rewritten afterward, ", True),
      ("because an earnings announcement can cause a large, unpredictable jump.")]},
    {"runs": [("Re-entered only if still wanted. ", True), ("If a stock is sold through an option, it "
      "is bought back only when it still ranks well, never chased above the strike.")]},
], Inches(0.7), Inches(y + 3.75), Inches(11.93), Inches(1.2), size=13.5, gap=8)
footer(s)

content("5 · The income overlay", "An honest limit on coverage", [
    {"runs": [("Options trade in blocks of 100 shares. ", True), ("A stock position must be at least "
      "100 shares before we can sell a call against it.")]},
    {"runs": [("Some positions are too small to cover. ", True), ("On higher-priced stocks, a normal "
      "position can fall short of 100 shares, so part of the book earns no premium.")]},
    {"runs": [("This is a known constraint, not a flaw in the idea. ", True), ("Position sizing and the "
      "investable amount determine how much of the book the overlay can reach, and that is managed "
      "deliberately.")]},
], subtitle="The overlay covers most, but not all, of the portfolio")

# =========================================================================== #
# § 6 — THE DATA PIPELINE
# =========================================================================== #
divider(6, "The data pipeline", "Where the numbers come from, and how they stay trustworthy")

s = slide(); y = header(s, "6 · The data pipeline", "The same sequence runs every trading day",
                        "Fully automated, and checked against the broker before any order is placed")
stages = [("Gather data", "prices and\nfinancials"), ("Score stocks", "the four\nfactors"),
          ("Build portfolio", "weights and\nlimits"), ("Add overlay", "covered\ncalls"),
          ("Risk checks", "pre-trade\nreview"), ("Place orders", "send to\nbroker")]
bx, by, bw, bh = 0.7, 2.7, 1.72, 1.35
agap = 0.32
for i, (lab, sub) in enumerate(stages):
    x = bx + i * (bw + agap)
    fill = ACCENT_DK if lab == "Add overlay" else INK
    rect(s, Inches(x), Inches(by), Inches(bw), Inches(bh), fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    text(s, Inches(x), Inches(by + 0.24), Inches(bw), Inches(0.5), lab, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, Inches(x + 0.05), Inches(by + 0.72), Inches(bw - 0.1), Inches(0.5), sub, size=9.5,
         color=RGBColor(0xB4, 0xC2, 0xCE), align=PP_ALIGN.CENTER, line_spacing=1.0)
    if i < len(stages) - 1:
        arrow(s, Inches(x + bw + 0.02), Inches(by + bh / 2 - 0.16), Inches(0.28), Inches(0.32))
text(s, Inches(0.7), Inches(4.75), Inches(11.93), Inches(0.8),
     "The portfolio is rebalanced once a month. The rest of the time the system monitors positions, "
     "manages the options, and keeps its records aligned with the broker.", size=14, color=TEXT, line_spacing=1.25)
footer(s)

s = slide(); y = header(s, "6 · The data pipeline", "Three data sources, each with a clear job",
                        "A small, low-cost, auditable set of inputs feeding the model")
image_fit(s, IMG_SOURCES, Inches(0.7), Inches(y + 0.05), Inches(11.93), Inches(4.75), valign="top")
footer(s)

s = slide(); y = header(s, "6 · The data pipeline", "Keeping the data trustworthy",
                        "The unglamorous half of the work, and where many strategies quietly fail")
image_fit(s, IMG_INTEGRITY, Inches(0.7), Inches(y + 0.1), Inches(11.93), Inches(4.6), valign="top")
footer(s)

# =========================================================================== #
# § 7 — EXECUTION, RISK & RESULTS
# =========================================================================== #
divider(7, "Execution, risk and results", "How orders are placed, how risk is controlled, and how it tested")

# Execution table
s = slide(); y = header(s, "7 · Execution and risk", "How orders reach the market",
                        "Order type is chosen to balance certainty of filling against control of price")
headers = ["Situation", "Order type", "Why"]
rows = [
    ["Large, highly liquid stock", "Market order", "Fills immediately, with negligible price impact"],
    ["Less liquid stock", "Marketable limit", "Fills the same session while bounding the price paid"],
    ["Selling a covered call", "Limit at mid-price", "Takes a fair premium without chasing the market"],
    ["Closing a call near expiry", "Market order", "Prioritises certainty of closing the position"],
]
cw = [3.2, 2.5, 6.23]
nrows = len(rows) + 1
gfx = s.shapes.add_table(nrows, 3, Inches(0.7), Inches(y + 0.05), Inches(sum(cw)), Inches(0.55 * nrows))
tbl = gfx.table; tbl.first_row = False; tbl.horz_banding = False
for i, c in enumerate(cw):
    tbl.columns[i].width = Inches(c)
for c, h_ in enumerate(headers):
    cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = INK; cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(cell, m, Inches(0.08))
    pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
    rn = pr.add_run(); rn.text = h_; rn.font.name = FONT; rn.font.size = Pt(12.5); rn.font.bold = True; rn.font.color.rgb = WHITE
for ri, row in enumerate(rows, 1):
    for c, val in enumerate(row):
        cell = tbl.cell(ri, c); cell.fill.solid(); cell.fill.fore_color.rgb = PANEL if ri % 2 == 0 else WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(cell, m, Inches(0.08))
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
        rn = pr.add_run(); rn.text = val; rn.font.name = FONT; rn.font.size = Pt(12); rn.font.bold = (c <= 1)
        rn.font.color.rgb = INK if c <= 1 else TEXT
bullets(s, [
    {"runs": [("Timing. ", True), ("The monthly rebalance runs mid-session, leaving a liquid window "
      "before the close for orders to fill.")]},
    {"runs": [("Unfilled orders. ", True), ("Cancelled after a short window and completed at the next "
      "run. Nothing is left resting overnight.")]},
], Inches(0.7), Inches(y + 0.05 + 0.55 * nrows + 0.25), Inches(11.93), Inches(1.2), size=13.5, gap=8)
footer(s)

s = slide(); y = header(s, "7 · Execution and risk", "Controlling risk",
                        "For a capital-preservation mandate, this is the most important part")
image_fit(s, IMG_RISK, Inches(0.7), Inches(y + 0.1), Inches(11.93), Inches(4.6), valign="top")
footer(s)

s = slide(); y = header(s, "7 · Execution and risk", "How the strategy was tested",
                        "A walk-forward simulation over 2021 to 2026, designed to resist over-fitting")
image_fit(s, IMG_METHOD, Inches(0.7), Inches(y + 0.35), Inches(11.93), Inches(3.0), valign="top")
text(s, Inches(0.7), Inches(y + 3.75), Inches(11.93), Inches(0.6),
     "Each month, the strategy runs on only the information available at the time, exactly as it "
     "would run live. The results, and their limitations, follow.", size=13.5, color=TEXT, line_spacing=1.2)
footer(s)

# Results
s = slide(); y = header(s, "7 · Results", "How the strategy performed in testing",
                        "The most dependable result is lower risk: smaller swings and shallower losses than the market")
image_fit(s, IMG_RR, Inches(0.55), Inches(y + 0.05), Inches(7.3), Inches(4.55), valign="top", halign="left")
cards = [("−16.7%", "Worst loss (drawdown) in the test,\nvs. −24% for the S&P 500", ACCENT_DK),
         ("15.5%", "Volatility, the size of the swings,\nvs. 16% for the S&P 500", INK),
         ("1.24", "Risk-adjusted return in the test,\nvs. 0.87 for the S&P 500", INK)]
cx = 8.15
for i, (big, lab, col) in enumerate(cards):
    cyy = y + 0.1 + i * 1.5
    rect(s, Inches(cx), Inches(cyy), Inches(4.48), Inches(1.32), fill=PANEL, line=PANEL_BD, line_w=1,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    rect(s, Inches(cx), Inches(cyy), Inches(0.09), Inches(1.32), fill=ACCENT)
    text(s, Inches(cx + 0.28), Inches(cyy), Inches(1.95), Inches(1.32), big, size=29, color=col,
         bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(cx + 2.1), Inches(cyy), Inches(2.25), Inches(1.32), lab, size=10.5, color=MUTED,
         line_spacing=1.05, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.7), Inches(6.5), Inches(11.93), Inches(0.55),
     [("Sharpe ratio", True), (" is return per unit of risk; ", False, MUTED),
      ("drawdown", True), (" is the largest peak-to-trough loss. Benchmarks are from index history; "
      "strategy figures from the simulation. Simulated and paper-traded, not live capital. "
      "Limitations on the next slide.", False, MUTED)], size=10.5, color=TEXT, line_spacing=1.12)
footer(s)

# ---- Limitations ---------------------------------------------------------- #
s = slide(); y = header(s, "7 · Results", "How to read these results",
                        "The results are encouraging, but they carry real limitations")
lim_l = [
    {"runs": [("Simulated, not live. ", True), ("A historical simulation and paper trading, not a track "
      "record, and not a prediction of future returns.")]},
    {"runs": [("One favourable period. ", True), ("The 2021 to 2026 window had specific conditions. A "
      "sustained bull market would likely show lower relative returns.")]},
    {"runs": [("Premium is partly modeled. ", True), ("Where real option prices are unavailable, the "
      "premium is estimated, so the overlay's contribution is uncertain.")]},
]
lim_r = [
    {"runs": [("Costs are estimates. ", True), ("Real trading costs, taxes, and slippage may run higher "
      "than modeled.")]},
    {"runs": [("Mild hindsight in the data. ", True), ("Some company financials are restated over time, "
      "making the factor results modestly optimistic.")]},
    {"runs": [("Figures are unlevered. ", True), ("Leverage would scale both returns and losses: larger "
      "gains, but larger drawdowns.")]},
]
bullets(s, lim_l, Inches(0.7), Inches(y + 0.2), Inches(5.7), Inches(4.3), size=14, gap=18)
rect(s, Inches(6.66), Inches(y + 0.1), Pt(1.2), Inches(4.3), fill=RULE)
bullets(s, lim_r, Inches(6.95), Inches(y + 0.2), Inches(5.7), Inches(4.3), size=14, gap=18)
text(s, Inches(0.7), Inches(6.5), Inches(11.93), Inches(0.5),
     [("The real verdict comes from live paper trading on real prices, not from more simulation.",
       True, INK)], size=13)
footer(s)

# =========================================================================== #
# § 8 — OPERATIONS
# =========================================================================== #
divider(8, "Operations", "A live, monitored system, running on its own")

s = slide(); y = header(s, "8 · Operations", "Monitoring: a live dashboard",
                        "One page, four tabs, updating itself from the live account")
image_fit(s, IMG_DASH, Inches(0.7), Inches(y + 0.1), Inches(11.93), Inches(4.4), valign="top")
text(s, Inches(0.7), Inches(y + 4.55), Inches(11.93), Inches(0.5),
     [("Overview · Portfolio · Performance · Backtest", True),
      ("    account value, holdings vs. target, factor tilt, growth vs. benchmarks, risk, and "
       "execution costs. Served read-only behind a password.", False, MUTED)], size=12)
footer(s)

content("8 · Operations", "Infrastructure and safeguards", [
    {"runs": [("Runs unattended in the cloud. ", True), ("Hosted on a dedicated server that runs the "
      "daily process and the dashboard on a fixed schedule.")]},
    {"runs": [("Restart-safe. ", True), ("Every component restarts automatically, with a watchdog that "
      "checks the system is alive.")]},
    {"runs": [("Backed up nightly. ", True), ("The operational database is copied to cloud storage each "
      "night, with a retention policy.")]},
    {"runs": [("Alerts by email. ", True), ("The system emails on the events that matter, such as a "
      "failed run or a risk check that blocks a trade.")]},
    {"runs": [("Shared securely. ", True), ("The dashboard is viewable read-only behind a password.")]},
], subtitle="Reliable, low-maintenance plumbing")

# =========================================================================== #
# § 9 — STATUS & OUTLOOK
# =========================================================================== #
divider(9, "Status and outlook", "Where the build stands, and what remains before live capital")

content("9 · Status and outlook", "Where the build stands today", [
    {"runs": [("The full strategy is built and tested. ", True), ("Data, factor scoring, portfolio "
      "construction, the options overlay, execution, and risk controls are all in place.")]},
    {"runs": [("It is running on a paper account. ", True), ("The system trades a simulated account "
      "with live market data, monitored by the dashboard.")]},
    {"runs": [("The remaining step is live-paper verification. ", True, ACCENT_DK),
      ("Confirm the full strategy fills, prices, and behaves correctly on the real account before any "
       "discussion of live capital.", False, INK)]},
], subtitle="Built, tested, and paper-trading, with one verification step ahead")

two_col("9 · Status and outlook", "Risks and what comes next",
        "Risks we monitor",
        [{"runs": [("Crowding in factors", True)]},
         {"lvl": 1, "runs": [("tracked by per-factor performance")]},
         {"runs": [("Capped upside in strong rallies", True)]},
         {"lvl": 1, "runs": [("the deliberate cost of the overlay")]},
         {"runs": [("Data-source reliability", True)]},
         {"lvl": 1, "runs": [("clean universe, reconciled each run")]},
         {"runs": [("The gap from test to live", True)]},
         {"lvl": 1, "runs": [("what paper trading is there to close")]}],
        "What comes next",
        [{"runs": [("Live-paper verification", True)]},
         {"lvl": 1, "runs": [("the immediate milestone")]},
         {"runs": [("Refine factor weights", True)]},
         {"lvl": 1, "runs": [("after more live evidence")]},
         {"runs": [("Broaden the income overlay", True)]},
         {"lvl": 1, "runs": [("more coverage, additional option strategies")]},
         {"runs": [("Define the go-live criteria", True)]},
         {"lvl": 1, "runs": [("before any real capital is committed")]}])

# ---- Summary -------------------------------------------------------------- #
s = slide(); rect(s, 0, 0, EMU_W, EMU_H, fill=INK); rect(s, 0, 0, Inches(0.22), EMU_H, fill=ACCENT)
glyph(s, 12.0, 0.85, d=0.15, gap=0.26)
text(s, Inches(0.9), Inches(0.85), Inches(11), Inches(0.4), "IN SUMMARY", size=13, color=ACCENT, bold=True)
summary = [
    ("A disciplined factor portfolio", "Four proven traits, computed transparently and weighted equally."),
    ("An income overlay that lowers risk", "Covered calls harvest a persistent premium and smooth the ride."),
    ("Built for capital preservation", "Firm diversification limits, a pre-trade risk check, and an emergency stop."),
    ("Tested, and honest about its limits", "Strong risk-adjusted returns in a rigorous simulation, with the assumptions stated."),
    ("Live and monitored", "Running on a paper account, one verification step from the go-live discussion."),
]
for i, (h_, d_) in enumerate(summary):
    yy = 1.55 + i * 1.02
    rect(s, Inches(0.9), Inches(yy + 0.05), Inches(0.13), Inches(0.7), fill=ACCENT)
    text(s, Inches(1.25), Inches(yy), Inches(11.2), Inches(0.4), h_, size=18, color=WHITE, bold=True)
    text(s, Inches(1.25), Inches(yy + 0.42), Inches(11.2), Inches(0.55), d_, size=13, color=RGBColor(0xA9, 0xB7, 0xC4), line_spacing=1.05)
footer(s)

# =========================================================================== #
# APPENDIX
# =========================================================================== #
divider(10, "Appendix", "Formulas, sources, detailed results, and a glossary")

# Formula reference
s = slide(); y = header(s, "Appendix · A1", "Factor formulas")
fcards = [("Value", F_VALUE), ("Momentum", F_MOM), ("Quality", F_QUAL), ("Low volatility", F_LOWVOL)]
for i, (lab, img) in enumerate(fcards):
    col = i % 2; row = i // 2
    formula_card(s, 0.7 + col * 6.13, y + 0.15 + row * 1.95, 5.7, 1.7, img, label=lab)
text(s, Inches(0.7), Inches(y + 4.1), Inches(11.93), Inches(0.5),
     [("z(·)", True, INK), (" is the standardising score from slide 8; each factor is combined "
      "equally into the final score.", False, MUTED)], size=12)
footer(s)

# Citations
s = slide(); y = header(s, "Appendix · A2", "Selected research")
cites = [
    ("Value", ["Fama & French (1992), The Cross-Section of Expected Stock Returns",
               "Fama & French (1993), Common Risk Factors in Stocks and Bonds",
               "Basu (1977), Investment Performance in Relation to P/E Ratios"]),
    ("Momentum", ["Jegadeesh & Titman (1993), Returns to Buying Winners and Selling Losers",
                  "Carhart (1997), On Persistence in Mutual Fund Performance",
                  "Asness, Moskowitz & Pedersen (2013), Value and Momentum Everywhere"]),
    ("Quality", ["Novy-Marx (2013), The Gross Profitability Premium",
                 "Asness, Frazzini & Pedersen (2019), Quality Minus Junk",
                 "Piotroski (2000), Value Investing and the F-Score"]),
    ("Low volatility", ["Ang, Hodrick, Xing & Zhang (2006), The Cross-Section of Volatility",
                        "Frazzini & Pedersen (2014), Betting Against Beta",
                        "Baker, Bradley & Wurgler (2011), Benchmarks as Limits to Arbitrage"]),
    ("Options premium", ["Whaley (2002), Return and Risk of the CBOE BuyWrite Index",
                         "Carr & Wu (2009), Variance Risk Premiums",
                         "Bakshi & Kapadia (2003), Delta-Hedged Gains and the Volatility Premium"]),
]
top = y + 0.05
for i, (cat, refs) in enumerate(cites):
    yy = top + i * 1.04
    text(s, Inches(0.7), Inches(yy), Inches(2.4), Inches(0.4), cat, size=13, color=ACCENT_DK, bold=True)
    for j, r in enumerate(refs):
        text(s, Inches(3.0), Inches(yy + j * 0.3), Inches(9.6), Inches(0.32), "·  " + r, size=11, color=TEXT)
footer(s)

# Detailed results
table_slide("Appendix · A3", "Detailed results, 2021 to 2026",
            ["", "Annual return", "Volatility", "Sharpe", "Max drawdown", "Premium / yr"],
            [["S&P 500", "+13.3%", "16.0%", "0.87", "−23.9%", "—"],
             ["S&P 500 BuyWrite (BXMD)", "+10.2%", "12.0%", "0.87", "−21.6%", "—"],
             ["Strategy, stocks only", "+10.9%", "~20%", "0.62", "−26%", "—"],
             ["Strategy + covered calls", "+19.6%", "15.5%", "1.24", "−16.7%", "~25%"],
             ["Strategy + put-wheel (optional)", "+18.7%", "12.9%", "1.40", "−12.8%", "—"]],
            [4.6, 1.6, 1.5, 1.15, 1.55, 1.53], hl={3}, fsize=11.5, hsize=10.5,
            subtitle="Strategy figures from the walk-forward simulation; benchmarks from index history",
            note="Benchmarks (S&P 500 and BXMD) computed from total-return index history over Sep 2021 to "
                 "Jun 2026, monthly. Strategy figures are from the walk-forward backtest on a 1x (unlevered) "
                 "basis. All results are simulated or paper-traded, not live capital.")

# VRP detail
stat_slide("Appendix · A4", "The volatility premium, measured",
           [("+5.2 pts", "Implied volatility sold (38.4%)\nover realized (33.2%)"),
            ("78%", "of months with implied\nabove realized volatility"),
            ("1.16x", "average implied-to-realized\nvolatility ratio"),
            ("22%", "of holdings called away\n(within the 30% target)")],
           subtitle="The premium the overlay collects is real and persistent",
           caption="Most of the premium compensates the risk of occasional large single-stock moves. The "
                   "30-delta strike is chosen to cap exactly that risk, which is why the overlay improves "
                   "returns even after accounting for it.")

# Parameters
table_slide("Appendix · A5", "Key parameters",
            ["Setting", "Value", "Setting", "Value"],
            [["Universe", "liquid US stocks", "Max per stock", "5%"],
             ["Holdings", "about 20", "Max per sector", "30%"],
             ["Factor weights", "25% each", "Invested / cash", "95% / 5%"],
             ["Momentum window", "12 months less 1", "Volatility window", "12 months"],
             ["Call delta target", "0.30", "Call expiry", "30 to 45 days"],
             ["Rebalance", "monthly", "Benchmarks", "S&P 500, BXMD"]],
            [3.1, 3.0, 3.0, 2.83], fsize=12, hsize=11,
            subtitle="Every setting lives in one configuration file, not scattered through the code")

# Glossary
s = slide(); y = header(s, "Appendix · A6", "Glossary")
glossary = [
    ("Factor", "A company trait linked to higher long-run returns"),
    ("z-score", "How far a value sits from the average, in standard deviations"),
    ("P/E, P/B", "Price vs. earnings, and price vs. net asset value"),
    ("ROE", "Profit as a percentage of shareholders' money"),
    ("Volatility", "How much a price moves around, a measure of risk"),
    ("Call option", "The right to buy a stock at a set price by a set date"),
    ("Premium", "The cash received for selling an option"),
    ("Strike", "The fixed price at which an option can be exercised"),
    ("Delta", "An option's rough chance of being exercised"),
    ("Implied volatility", "The price movement the market expects, priced into options"),
    ("Assignment", "Being required to sell the stock when a call is exercised"),
    ("Sharpe ratio", "Return earned per unit of risk taken"),
]
top = y + 0.1
for i, (term, defn) in enumerate(glossary):
    col = i // 6; row = i % 6
    xx = 0.7 + col * 6.15; yy = top + row * 0.78
    text(s, Inches(xx), Inches(yy), Inches(5.9), Inches(0.32), term, size=13, color=INK, bold=True)
    text(s, Inches(xx), Inches(yy + 0.3), Inches(5.9), Inches(0.45), defn, size=11.5, color=MUTED, line_spacing=1.0)
footer(s)

# --------------------------------------------------------------------------- #
out = os.path.join(HERE, "Systematic-Factor-Income-Fund.pptx")
prs.save(out)
print(f"saved {out}  —  {len(prs.slides._sldIdLst)} slides")
